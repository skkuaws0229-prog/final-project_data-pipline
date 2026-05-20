from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("docs/aws_icon_architecture_drug_repurposing_v1.pptx")
ICON_ROOT = Path("/private/tmp/aws_icons_04302026/Architecture-Service-Icons_04302026")
GROUP_ROOT = Path("/private/tmp/aws_icons_04302026/Architecture-Group-Icons_04302026")


COLORS = {
    "ink": RGBColor(17, 24, 39),
    "muted": RGBColor(75, 85, 99),
    "line": RGBColor(55, 65, 81),
    "aws_orange": RGBColor(255, 153, 0),
    "aws_blue": RGBColor(35, 47, 62),
    "bg": RGBColor(247, 249, 252),
    "box": RGBColor(255, 255, 255),
    "dev": RGBColor(236, 253, 245),
    "runtime": RGBColor(239, 246, 255),
    "data": RGBColor(255, 247, 237),
    "future": RGBColor(250, 245, 255),
    "handoff": RGBColor(248, 250, 252),
}


ICONS = {
    "AWS Cloud": GROUP_ROOT / "AWS-Cloud-logo_32.png",
    "S3": ICON_ROOT / "Arch_Storage/64/Arch_Amazon-Simple-Storage-Service_64.png",
    "EC2": ICON_ROOT / "Arch_Compute/64/Arch_Amazon-EC2_64.png",
    "Lambda": ICON_ROOT / "Arch_Compute/64/Arch_AWS-Lambda_64.png",
    "Step Functions": ICON_ROOT / "Arch_Application-Integration/64/Arch_AWS-Step-Functions_64.png",
    "RDS": ICON_ROOT / "Arch_Databases/64/Arch_Amazon-RDS_64.png",
    "OpenSearch": ICON_ROOT / "Arch_Analytics/64/Arch_Amazon-OpenSearch-Service_64.png",
    "SageMaker": ICON_ROOT / "Arch_Artificial-Intelligence/64/Arch_Amazon-SageMaker-AI_64.png",
    "Bedrock": ICON_ROOT / "Arch_Artificial-Intelligence/64/Arch_Amazon-Bedrock_64.png",
    "API Gateway": ICON_ROOT / "Arch_Networking-Content-Delivery/64/Arch_Amazon-API-Gateway_64.png",
    "CloudWatch": ICON_ROOT / "Arch_Management-Tools/64/Arch_Amazon-CloudWatch_64.png",
    "Secrets": ICON_ROOT / "Arch_Security-Identity/64/Arch_AWS-Secrets-Manager_64.png",
}


def add_text(slide, text, x, y, w, h, size=12, bold=False, color="ink", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color] if isinstance(color, str) else color
    return box


def add_zone(slide, title, x, y, w, h, fill, line=RGBColor(148, 163, 184)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    add_text(slide, title, x + 0.12, y + 0.05, w - 0.24, 0.28, size=10, bold=True, color="ink")
    return shape


def add_service(slide, label, icon_key, x, y, w=1.25, h=0.85, subtitle=None):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["box"]
    card.line.color.rgb = RGBColor(203, 213, 225)
    card.line.width = Pt(1)
    icon_path = ICONS.get(icon_key)
    if icon_path and icon_path.exists():
        slide.shapes.add_picture(str(icon_path), Inches(x + 0.08), Inches(y + 0.15), width=Inches(0.33), height=Inches(0.33))
    add_text(slide, label, x + 0.44, y + 0.10, w - 0.50, 0.34, size=8.3, bold=True, color="ink")
    if subtitle:
        add_text(slide, subtitle, x + 0.10, y + 0.48, w - 0.20, 0.24, size=6.9, color="muted")
    return card


def add_external(slide, label, x, y, w=1.15, h=0.62):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(148, 163, 184)
    shape.line.width = Pt(1)
    add_text(slide, label, x + 0.06, y + 0.08, w - 0.12, h - 0.14, size=8, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2, label=None, dashed=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS["line"]
    line.line.width = Pt(1.2)
    line.line.end_arrowhead = True
    if dashed:
        line.line.dash_style = 4
    if label:
        lx = min(x1, x2) + abs(x2 - x1) / 2 - 0.35
        ly = min(y1, y2) + abs(y2 - y1) / 2 - 0.16
        add_text(slide, label, lx, ly, 0.78, 0.22, size=6.4, color="muted", align=PP_ALIGN.CENTER)
    return line


def add_step(slide, n, x, y):
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.30), Inches(0.30))
    circ.fill.solid()
    circ.fill.fore_color.rgb = COLORS["aws_orange"]
    circ.line.color.rgb = RGBColor(234, 88, 12)
    circ.line.width = Pt(0.5)
    add_text(slide, str(n), x + 0.02, y + 0.035, 0.26, 0.20, size=7.5, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    return circ


def add_person(slide, label, x, y):
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.22), Inches(y), Inches(0.22), Inches(0.22))
    head.fill.solid()
    head.fill.fore_color.rgb = RGBColor(255, 255, 255)
    head.line.color.rgb = COLORS["line"]
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(x + 0.08), Inches(y + 0.20), Inches(0.52), Inches(0.40))
    body.line.color.rgb = COLORS["line"]
    body.line.width = Pt(1)
    add_text(slide, label, x, y + 0.52, 0.72, 0.26, size=7.2, align=PP_ALIGN.CENTER)


def slide_main(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["bg"]

    add_text(slide, "Precision Drug Repurposing Platform - AWS Architecture", 0.35, 0.14, 8.7, 0.28, size=15, bold=True)
    add_text(slide, "현재 로컬/Docker 검증 구조 + AWS 배포/자동화 확장 구조", 0.35, 0.46, 6.2, 0.20, size=8.5, color="muted")
    if ICONS["AWS Cloud"].exists():
        slide.shapes.add_picture(str(ICONS["AWS Cloud"]), Inches(12.15), Inches(0.17), width=Inches(0.30), height=Inches(0.30))
    add_text(slide, "AWS icon set 2026-04-30", 10.95, 0.46, 1.65, 0.18, size=6.5, color="muted", align=PP_ALIGN.RIGHT)

    add_person(slide, "Researcher", 0.20, 1.45)
    add_person(slide, "Frontend Dev", 0.16, 4.95)

    add_zone(slide, "Frontend / Team Client", 1.05, 0.85, 1.65, 5.90, COLORS["handoff"])
    react = add_external(slide, "React Vite\n5174", 1.28, 1.55)
    browser = add_external(slide, "Browser QA\nCandidates / KG / 3D", 1.28, 2.68, h=0.78)
    chatbot = add_external(slide, "Chatbot UI\nfallback + team API", 1.28, 4.02, h=0.78)

    add_zone(slide, "Drug Service AWS Account / EC2 Runtime", 2.95, 0.85, 5.25, 5.90, COLORS["runtime"])
    ec2 = add_service(slide, "Amazon EC2", "EC2", 3.25, 1.25, subtitle="Docker Compose host")
    api = add_service(slide, "FastAPI", "API Gateway", 4.95, 1.25, subtitle="REST API layer")
    pg = add_service(slide, "PostgreSQL", "RDS", 3.25, 2.45, subtitle="candidate / final / structures")
    neo = add_service(slide, "Neo4j KG", "EC2", 4.95, 2.45, subtitle="disease-drug-target graph")
    os = add_service(slide, "OpenSearch", "OpenSearch", 6.65, 2.45, subtitle="search / RAG evidence")
    af = add_service(slide, "AlphaFold Proxy", "S3", 3.25, 3.75, subtitle="CIF file endpoint")
    kg = add_service(slide, "KG UI API", "API Gateway", 4.95, 3.75, subtitle="ui-basic / summary")
    explain = add_service(slide, "Explanation Context", "Bedrock", 6.65, 3.75, subtitle="Bedrock-ready JSON")
    assistant = add_service(slide, "Assistant API", "Bedrock", 4.95, 5.05, subtitle="suggested Q / ask mock")

    add_zone(slide, "Data Lake / Evidence Store", 8.45, 0.85, 2.00, 5.90, COLORS["data"])
    s3raw = add_service(slide, "Amazon S3", "S3", 8.78, 1.25, subtitle="raw / results / structures")
    s3pkg = add_service(slide, "S3 Handoff", "S3", 8.78, 2.45, subtitle="docs / API contracts")
    gh = add_external(slide, "GitHub\nfinal-project_data-pipline", 8.78, 3.72, 1.25, 0.78)
    gcs = add_external(slide, "GCS backup\nfuture option", 8.78, 5.05, 1.25, 0.70)

    add_zone(slide, "Future Automation Layer", 10.70, 0.85, 2.15, 5.90, COLORS["future"])
    lam = add_service(slide, "AWS Lambda", "Lambda", 11.02, 1.25, subtitle="S3 event trigger")
    sf = add_service(slide, "Step Functions", "Step Functions", 11.02, 2.45, subtitle="pipeline orchestration")
    sm = add_service(slide, "SageMaker AI", "SageMaker", 11.02, 3.65, subtitle="heavy pipeline jobs")
    br = add_service(slide, "Amazon Bedrock", "Bedrock", 11.02, 4.85, subtitle="RAG / LLM chatbot")

    add_arrow(slide, 0.88, 1.80, 1.28, 1.80)
    add_arrow(slide, 2.43, 1.85, 4.95, 1.63, "API call")
    add_arrow(slide, 5.58, 1.98, 3.86, 2.45)
    add_arrow(slide, 5.65, 1.98, 5.55, 2.45)
    add_arrow(slide, 5.75, 1.98, 7.18, 2.45)
    add_arrow(slide, 3.88, 3.75, 3.88, 3.30)
    add_arrow(slide, 5.55, 3.75, 5.55, 3.30)
    add_arrow(slide, 7.22, 3.75, 7.22, 3.30)
    add_arrow(slide, 8.78, 1.65, 8.02, 1.65, "load")
    add_arrow(slide, 8.78, 2.85, 7.95, 2.85, "evidence")
    add_arrow(slide, 10.03, 1.65, 11.02, 1.65, "future")
    add_arrow(slide, 11.62, 2.10, 11.62, 2.45)
    add_arrow(slide, 11.62, 3.30, 11.62, 3.65)
    add_arrow(slide, 11.62, 4.50, 11.62, 4.85)
    add_arrow(slide, 10.98, 5.20, 7.90, 4.22, dashed=True)
    add_arrow(slide, 9.42, 4.50, 9.42, 5.05, dashed=True)
    add_arrow(slide, 2.43, 4.42, 4.95, 5.35, "assistant")

    for i, xy in enumerate([(0.88, 1.58), (2.72, 1.55), (4.62, 1.55), (3.02, 2.74), (4.72, 2.74), (6.42, 2.74), (3.02, 4.08), (4.72, 4.08), (8.50, 1.56), (10.72, 1.56), (10.72, 2.76), (10.72, 4.00)], 1):
        add_step(slide, i, *xy)

    add_text(slide, "검은 실선: 현재 연결 완료 / 점선: 향후 자동화·Bedrock·GCS 확장", 0.35, 6.92, 6.8, 0.22, size=8, color="muted")
    add_text(slide, "No TxGNN runtime: 비용/매칭 이슈로 제외, Neo4j path scoring + KG embedding + RAG 설명 중심", 7.1, 6.92, 5.6, 0.22, size=8, color="muted", align=PP_ALIGN.RIGHT)


def slide_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 255, 255)
    add_text(slide, "Backend-to-Frontend Workflow", 0.35, 0.18, 5.5, 0.32, size=15, bold=True)
    add_text(slide, "프론트 QA 기준으로 실제 연결된 API와 향후 Bedrock/RAG 진입점을 정리", 0.35, 0.52, 7.0, 0.22, size=8.5, color="muted")

    cols = [
        ("1. Data Ingestion", "S3 + local CSV\n11 diseases\ncandidate/final split", "S3", 0.55, COLORS["data"]),
        ("2. Storage", "PostgreSQL\nNeo4j KG\nOpenSearch", "RDS", 2.65, COLORS["runtime"]),
        ("3. API Layer", "FastAPI\nCandidates / Graph\nStructures / Search", "API Gateway", 4.75, COLORS["handoff"]),
        ("4. Frontend QA", "React 5174\nCandidates 15 vs 60\nKG + AlphaFold", "CloudWatch", 6.85, COLORS["dev"]),
        ("5. Explanation", "Explanation context\nAssistant API mock\nBedrock-ready", "Bedrock", 8.95, COLORS["future"]),
        ("6. Automation Later", "Lambda\nStep Functions\nSageMaker jobs", "Step Functions", 11.05, COLORS["future"]),
    ]
    prev_x = None
    for idx, (title, body, icon, x, fill) in enumerate(cols, 1):
        add_zone(slide, title, x, 1.25, 1.55, 4.85, fill)
        add_service(slide, title.split(". ", 1)[1], icon, x + 0.18, 1.85, 1.18, 0.95)
        add_text(slide, body, x + 0.16, 3.05, 1.22, 1.35, size=8, color="ink", align=PP_ALIGN.CENTER)
        add_step(slide, idx, x + 0.60, 5.36)
        if prev_x is not None:
            add_arrow(slide, prev_x + 1.55, 3.15, x, 3.15)
        prev_x = x

    add_text(slide, "프론트 전달 핵심 API", 0.55, 6.35, 2.2, 0.24, size=10, bold=True)
    api_list = (
        "/v1/diseases/{code}/final-candidates\n"
        "/api/diseases/{code}/candidates\n"
        "/api/graph/{code}/ui-basic\n"
        "/api/structures/targets?q=JAK\n"
        "/api/explanation-context\n"
        "/api/assistant/{code}/ask"
    )
    add_text(slide, api_list, 2.00, 6.15, 4.25, 0.72, size=7.1, color="muted")
    add_text(slide, "산출물: PPTX는 발표/수정용, SVG는 문서/README 삽입용, draw.io는 보조 편집용", 7.0, 6.35, 5.5, 0.24, size=8, color="muted", align=PP_ALIGN.RIGHT)


def main():
    missing = [str(p) for p in ICONS.values() if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing AWS icon files:\n" + "\n".join(missing))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_main(prs)
    slide_workflow(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT.resolve())


if __name__ == "__main__":
    main()
